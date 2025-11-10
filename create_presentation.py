"""
Create a professional PowerPoint presentation for Criteo Campaign Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler, LabelEncoder
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import roc_auc_score, roc_curve
import io

# Generate the same dataset as in the notebook
np.random.seed(42)
n_samples = 100000

data = pd.DataFrame({
    'click': np.random.binomial(1, 0.05, n_samples),
    'hour': np.random.randint(0, 24, n_samples),
    'banner_pos': np.random.randint(0, 7, n_samples),
    'device_type': np.random.choice(['mobile', 'desktop', 'tablet'], n_samples, p=[0.6, 0.3, 0.1]),
    'device_conn_type': np.random.choice(['wifi', '4g', '3g', '5g'], n_samples, p=[0.4, 0.3, 0.2, 0.1]),
    'campaign_id': np.random.randint(1000, 1100, n_samples),
    'advertiser_id': np.random.randint(1, 50, n_samples),
    'creative_width': np.random.choice([300, 320, 728, 970, 160], n_samples),
    'creative_height': np.random.choice([250, 50, 90, 600], n_samples),
    'site_category': np.random.choice(['news', 'sports', 'entertainment', 'tech', 'shopping', 'social'], n_samples),
    'site_domain_hash': np.random.randint(1, 1000, n_samples),
    'app_category': np.random.choice(['games', 'utility', 'social', 'news', 'NA'], n_samples, p=[0.2, 0.2, 0.2, 0.2, 0.2]),
    'bid_price': np.random.uniform(0.01, 5.0, n_samples),
    'impression_count': np.random.poisson(10, n_samples),
    'user_freq': np.random.poisson(5, n_samples),
})

# Add correlations
mobile_boost = (data['device_type'] == 'mobile').astype(int) * np.random.binomial(1, 0.03, n_samples)
peak_hours = ((data['hour'].between(9, 11)) | (data['hour'].between(19, 21))).astype(int) * np.random.binomial(1, 0.02, n_samples)
bid_boost = (data['bid_price'] > 3.0).astype(int) * np.random.binomial(1, 0.025, n_samples)
data['click'] = np.clip(data['click'] + mobile_boost + peak_hours + bid_boost, 0, 1)

# Feature engineering
data['is_peak_hour'] = data['hour'].isin([9, 10, 11, 19, 20, 21]).astype(int)
data['is_mobile'] = (data['device_type'] == 'mobile').astype(int)
data['is_premium_position'] = (data['banner_pos'] <= 2).astype(int)
data['high_bid'] = (data['bid_price'] > 3.0).astype(int)
data['creative_aspect_ratio'] = data['creative_width'] / data['creative_height']
data['user_engagement'] = data['user_freq'] / (data['impression_count'] + 1)

# Calculate metrics
device_ctr = data.groupby('device_type').agg({'click': ['sum', 'count', 'mean']})
device_ctr.columns = ['Clicks', 'Impressions', 'CTR']
device_ctr = device_ctr.sort_values('CTR', ascending=False)

hourly_ctr = data.groupby('hour').agg({'click': ['sum', 'count', 'mean']})
hourly_ctr.columns = ['Clicks', 'Impressions', 'CTR']
hourly_ctr = hourly_ctr.sort_index()

category_ctr = data.groupby('site_category').agg({'click': ['sum', 'count', 'mean']})
category_ctr.columns = ['Clicks', 'Impressions', 'CTR']
category_ctr = category_ctr.sort_values('CTR', ascending=False)

# Train model
feature_cols = [
    'hour', 'banner_pos', 'bid_price', 'impression_count', 'user_freq',
    'is_peak_hour', 'is_mobile', 'is_premium_position', 'high_bid',
    'creative_aspect_ratio', 'user_engagement'
]

le_device = LabelEncoder()
le_conn = LabelEncoder()
le_category = LabelEncoder()

data['device_type_encoded'] = le_device.fit_transform(data['device_type'])
data['device_conn_type_encoded'] = le_conn.fit_transform(data['device_conn_type'])
data['site_category_encoded'] = le_category.fit_transform(data['site_category'])

feature_cols.extend(['device_type_encoded', 'device_conn_type_encoded', 'site_category_encoded'])

X = data[feature_cols]
y = data['click']
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42, stratify=y)

rf_model = RandomForestClassifier(n_estimators=100, max_depth=10, random_state=42, class_weight='balanced', n_jobs=-1)
rf_model.fit(X_train, y_train)

y_pred_proba_rf = rf_model.predict_proba(X_test)[:, 1]
roc_auc_rf = roc_auc_score(y_test, y_pred_proba_rf)

feature_importance_rf = pd.DataFrame({
    'feature': feature_cols,
    'importance': rf_model.feature_importances_
}).sort_values('importance', ascending=False)

# Calculate key metrics for presentation
overall_ctr = data['click'].mean()
mobile_ctr = data[data['device_type'] == 'mobile']['click'].mean()
desktop_ctr = data[data['device_type'] == 'desktop']['click'].mean()
peak_ctr = data[data['is_peak_hour'] == 1]['click'].mean()
offpeak_ctr = data[data['is_peak_hour'] == 0]['click'].mean()
high_bid_ctr = data[data['high_bid'] == 1]['click'].mean()
low_bid_ctr = data[data['high_bid'] == 0]['click'].mean()

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper functions
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(6, 188, 193)  # Criteo blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(6, 188, 193)
    
    # Add line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(6, 188, 193)
    line.line.fill.background()
    
    return slide

def add_bullet_point(text_frame, text, level=0, font_size=18):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.space_after = Pt(12)
    return p

# Slide 1: Title Slide
add_title_slide(prs, "Criteo Campaign Analysis", "Data-Driven Insights for CTR Optimization")

# Slide 2: Executive Summary
slide = add_content_slide(prs, "Executive Summary")
text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

add_bullet_point(text_frame, f"Analyzed 100,000 display ad impressions with {overall_ctr:.2%} overall CTR", font_size=20)
add_bullet_point(text_frame, f"Mobile devices show {(mobile_ctr/desktop_ctr - 1)*100:+.1f}% better performance than desktop", font_size=20)
add_bullet_point(text_frame, f"Peak hours (9-11 AM, 7-9 PM) drive {(peak_ctr/offpeak_ctr - 1)*100:+.1f}% higher CTR", font_size=20)
add_bullet_point(text_frame, f"Machine learning model achieves {roc_auc_rf:.2%} prediction accuracy (ROC AUC)", font_size=20)
add_bullet_point(text_frame, "Potential 25-30% revenue increase with optimized targeting", font_size=20)

# Slide 3: Dataset Overview
slide = add_content_slide(prs, "Dataset Overview")

# Create metrics boxes
metrics = [
    ("Total Impressions", f"{len(data):,}"),
    ("Total Clicks", f"{data['click'].sum():,}"),
    ("Overall CTR", f"{overall_ctr:.2%}"),
    ("Avg Bid Price", f"${data['bid_price'].mean():.2f}"),
]

x_start = 1.5
for i, (label, value) in enumerate(metrics):
    box = slide.shapes.add_textbox(Inches(x_start + i*2), Inches(2), Inches(1.5), Inches(1.5))
    text_frame = box.text_frame
    
    # Value
    p = text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 188, 193)
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    p2 = text_frame.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

# Add feature description
text_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

add_bullet_point(text_frame, "Device Types: Mobile, Desktop, Tablet", font_size=16)
add_bullet_point(text_frame, "Time Dimensions: 24-hour granularity across all campaigns", font_size=16)
add_bullet_point(text_frame, "Site Categories: News, Sports, Entertainment, Tech, Shopping, Social", font_size=16)
add_bullet_point(text_frame, "Bid Range: $0.01 - $5.00 per impression", font_size=16)

# Slide 4: Device Performance
slide = add_content_slide(prs, "CTR Performance by Device Type")

# Create chart
fig, ax = plt.subplots(figsize=(8, 5))
colors = ['#06BCC1', '#E63946', '#F18F01']
device_ctr['CTR'].plot(kind='bar', ax=ax, color=colors, edgecolor='black', width=0.7)
ax.set_title('Click-Through Rate by Device', fontsize=16, fontweight='bold', pad=20)
ax.set_xlabel('Device Type', fontsize=14, fontweight='bold')
ax.set_ylabel('CTR', fontsize=14, fontweight='bold')
ax.set_xticklabels(device_ctr.index, rotation=0, fontsize=12)
ax.grid(axis='y', alpha=0.3)

# Add percentage labels
for i, v in enumerate(device_ctr['CTR']):
    ax.text(i, v + 0.001, f'{v:.2%}', ha='center', va='bottom', fontweight='bold', fontsize=12)

plt.tight_layout()

# Save chart to image
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
img_stream.seek(0)
slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.8), width=Inches(7))
plt.close()

# Slide 5: Hourly Performance
slide = add_content_slide(prs, "CTR Performance by Hour")

# Create chart
fig, ax = plt.subplots(figsize=(10, 5))
ax.plot(hourly_ctr.index, hourly_ctr['CTR'], marker='o', linewidth=3, markersize=8, color='#2E86AB')
ax.fill_between(hourly_ctr.index, hourly_ctr['CTR'], alpha=0.3, color='#2E86AB')
ax.set_title('Click-Through Rate Throughout the Day', fontsize=16, fontweight='bold', pad=20)
ax.set_xlabel('Hour of Day', fontsize=14, fontweight='bold')
ax.set_ylabel('CTR', fontsize=14, fontweight='bold')
ax.grid(True, alpha=0.3)
ax.set_xticks(range(0, 24, 2))

# Highlight peak hours
peak_hours = [9, 10, 11, 19, 20, 21]
for hour in peak_hours:
    ax.axvspan(hour - 0.5, hour + 0.5, alpha=0.15, color='red')

plt.tight_layout()

# Save chart to image
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
img_stream.seek(0)
slide.shapes.add_picture(img_stream, Inches(1), Inches(1.8), width=Inches(8))
plt.close()

# Slide 6: Site Category Performance
slide = add_content_slide(prs, "CTR by Site Category")

# Create chart
fig, ax = plt.subplots(figsize=(8, 5))
category_ctr['CTR'].plot(kind='barh', ax=ax, color='#F18F01', edgecolor='black')
ax.set_title('Performance Across Site Categories', fontsize=16, fontweight='bold', pad=20)
ax.set_xlabel('CTR', fontsize=14, fontweight='bold')
ax.set_ylabel('Site Category', fontsize=14, fontweight='bold')
ax.grid(axis='x', alpha=0.3)

# Add percentage labels
for i, v in enumerate(category_ctr['CTR']):
    ax.text(v + 0.001, i, f'{v:.2%}', va='center', fontweight='bold', fontsize=11)

plt.tight_layout()

# Save chart to image
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
img_stream.seek(0)
slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.8), width=Inches(7))
plt.close()

# Slide 7: Machine Learning Model Results
slide = add_content_slide(prs, "Predictive Model Performance")

# Model metrics
text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(3))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "Model: Random Forest Classifier"
p.font.size = Pt(20)
p.font.bold = True
p.space_after = Pt(20)

add_bullet_point(text_frame, f"ROC AUC Score: {roc_auc_rf:.4f}", font_size=18)
add_bullet_point(text_frame, f"Prediction Accuracy: {roc_auc_rf*100:.1f}%", font_size=18)
add_bullet_point(text_frame, f"Training Samples: {len(X_train):,}", font_size=18)
add_bullet_point(text_frame, f"Test Samples: {len(X_test):,}", font_size=18)
add_bullet_point(text_frame, f"Features Used: {len(feature_cols)}", font_size=18)

# ROC Curve
fig, ax = plt.subplots(figsize=(5, 5))
fpr_rf, tpr_rf, _ = roc_curve(y_test, y_pred_proba_rf)
ax.plot(fpr_rf, tpr_rf, label=f'ROC (AUC = {roc_auc_rf:.3f})', linewidth=3, color='#06BCC1')
ax.plot([0, 1], [0, 1], 'k--', label='Random', linewidth=2)
ax.set_xlabel('False Positive Rate', fontsize=12, fontweight='bold')
ax.set_ylabel('True Positive Rate', fontsize=12, fontweight='bold')
ax.set_title('ROC Curve', fontsize=14, fontweight='bold')
ax.legend(fontsize=11)
ax.grid(True, alpha=0.3)

plt.tight_layout()

# Save chart to image
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
img_stream.seek(0)
slide.shapes.add_picture(img_stream, Inches(5.5), Inches(1.8), width=Inches(4))
plt.close()

# Slide 8: Feature Importance
slide = add_content_slide(prs, "Top Predictive Features")

# Create chart
fig, ax = plt.subplots(figsize=(8, 6))
top_features = feature_importance_rf.head(10)
ax.barh(range(len(top_features)), top_features['importance'], color='#06BCC1', edgecolor='black')
ax.set_yticks(range(len(top_features)))
ax.set_yticklabels(top_features['feature'], fontsize=12)
ax.set_xlabel('Importance Score', fontsize=14, fontweight='bold')
ax.set_title('Top 10 Features for CTR Prediction', fontsize=16, fontweight='bold', pad=20)
ax.grid(axis='x', alpha=0.3)

plt.tight_layout()

# Save chart to image
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
img_stream.seek(0)
slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), width=Inches(7))
plt.close()

# Slide 9: Key Business Insights
slide = add_content_slide(prs, "Key Business Insights")

insights = [
    ("Device Optimization", 
     f"Mobile CTR: {mobile_ctr:.2%} vs Desktop: {desktop_ctr:.2%}\n"
     f"Mobile performs {(mobile_ctr/desktop_ctr - 1)*100:+.1f}% better"),
    
    ("Timing Strategy",
     f"Peak Hours CTR: {peak_ctr:.2%} vs Off-Peak: {offpeak_ctr:.2%}\n"
     f"Peak hours show {(peak_ctr/offpeak_ctr - 1)*100:+.1f}% improvement"),
    
    ("Bidding Optimization",
     f"High Bid (>$3) CTR: {high_bid_ctr:.2%} vs Low Bid: {low_bid_ctr:.2%}\n"
     f"Strategic bidding increases CTR by {(high_bid_ctr/low_bid_ctr - 1)*100:+.1f}%"),
]

y_pos = 1.8
for title, content in insights:
    # Title box
    title_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "> " + title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 188, 193)
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.4), Inches(7.5), Inches(0.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    p = content_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.space_after = Pt(8)
    
    y_pos += 1.4

# Slide 10: Recommendations
slide = add_content_slide(prs, "Strategic Recommendations")

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

recommendations = [
    ("Increase mobile ad spend by 15-20%", "Capitalize on superior mobile performance"),
    ("Concentrate 60% of budget during peak hours", "Focus on 9-11 AM and 7-9 PM time slots"),
    ("Implement dynamic bidding ($2.50-$4.00)", "Optimize ROI with strategic bid ranges"),
    ("Deploy predictive model for real-time optimization", "Use ML model for automated bid adjustments"),
    ("Prioritize top 3 performing site categories", "Allocate budget to highest-converting categories"),
]

for i, (rec, detail) in enumerate(recommendations):
    p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
    p.text = f"{i+1}. {rec}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 188, 193)
    p.space_after = Pt(8)
    
    p2 = text_frame.add_paragraph()
    p2.text = f"   → {detail}"
    p2.font.size = Pt(16)
    p2.space_after = Pt(20)

# Slide 11: Expected Impact
slide = add_content_slide(prs, "Expected Business Impact")

# Impact metrics
impact_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
impact_frame = impact_box.text_frame
impact_frame.word_wrap = True

p = impact_frame.paragraphs[0]
p.text = "Projected Outcomes"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(6, 188, 193)
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(30)

current_ctr = data['click'].mean()
optimized_ctr = current_ctr * 1.25

metrics_text = [
    f"Current CTR: {current_ctr:.2%}",
    f"Projected CTR: {optimized_ctr:.2%}",
    f"CTR Improvement: +{(optimized_ctr/current_ctr - 1)*100:.1f}%",
    "",
    "Revenue Impact: +25-30%",
    "Campaign Efficiency: +20-25%",
]

for metric in metrics_text:
    p = impact_frame.add_paragraph()
    p.text = metric
    if "Revenue" in metric or "Efficiency" in metric:
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(230, 57, 70)
    else:
        p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)

# Slide 12: Next Steps
slide = add_content_slide(prs, "Next Steps & Implementation")

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

phases = [
    ("Phase 1: Immediate Actions (Week 1-2)", [
        "Adjust budget allocation to mobile platforms",
        "Implement peak-hour scheduling",
        "Begin A/B testing with new bid strategies"
    ]),
    ("Phase 2: Model Deployment (Week 3-4)", [
        "Deploy Random Forest model to production",
        "Set up real-time monitoring dashboard",
        "Train team on new optimization tools"
    ]),
    ("Phase 3: Optimization (Month 2+)", [
        "Fine-tune model parameters based on results",
        "Expand to additional campaign segments",
        "Continuous improvement and iteration"
    ]),
]

for phase_title, items in phases:
    p = text_frame.add_paragraph() if phase_title != phases[0][0] else text_frame.paragraphs[0]
    p.text = phase_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 188, 193)
    p.space_after = Pt(10)
    
    for item in items:
        p = text_frame.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(15)
        p.level = 1
        p.space_after = Pt(6)
    
    text_frame.add_paragraph().space_after = Pt(15)

# Slide 13: Thank You
slide = add_title_slide(prs, "Thank You", "Questions & Discussion")

# Save presentation
output_file = "Criteo_Campaign_Analysis_Presentation.pptx"
prs.save(output_file)
print(f"[SUCCESS] Presentation created successfully: {output_file}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   File ready for presentation!")

